from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt():
    prs = Presentation()

    # Slide content data
    slides_data = [
        {
            "title": "AI-Powered Coding Agent Framework",
            "content": [
                "What it is: A professional TypeScript-based agentic framework for AI-driven code generation.",
                "Core Purpose: To bridge the gap between LLM chat and actual software engineering by providing a terminal-based interface (CLI/TUI) that can autonomously read, write, and execute code.",
                "Primary Goal: Transform high-level natural language requirements into verified, working code implementations."
            ]
        },
        {
            "title": "Powerful Capabilities & Versatility",
            "content": [
                "Flexible Agent Modes: Supports Execute (planner-driven), Plan (manual approval), and ReAct (dynamic reasoning).",
                "Advanced TUI: Interactive terminal UI for real-time streaming of agent thoughts, task tracking, and status monitoring.",
                "Backend Agnostic: Seamless integration with OpenAI, Anthropic, and local models via Ollama.",
                "Domain Skill System: Pre-defined guidelines for specific frameworks (React, Flask, SQL) to improve code quality.",
                "Built-in Benchmarking: Integrated harness for performance evaluation using SWE-bench Lite and HumanEval."
            ]
        },
        {
            "title": "Core Orchestration: MasterCoordinator",
            "content": [
                "The MasterCoordinator: Acts as the central nervous system, managing task routing, state persistence, and sub-agent spawning.",
                "Task Lifecycle: Handles the transition from user prompt -> task creation -> routing -> execution.",
                "Concurrency Control: Implements FileLockManager to ensure multiple agents don't corrupt the same files simultaneously.",
                "Human-in-the-loop: Integrated clarification mechanism to pause and request user input before critical actions."
            ]
        },
        {
            "title": "The Planner-Executor Pattern",
            "content": [
                "The Planner: Analyzes the request and decomposes it into a structured sequence of PlanSteps (e.g., Tool Loop -> Code Change -> Verify).",
                "The Executor: Follows the plan strictly, utilizing a suite of tools: repo_map, edit_file/read_file, and bash.",
                "The Feedback Loop: An iterative process of Inspect -> Design -> Implement -> Verify, reducing hallucinations."
            ]
        },
        {
            "title": "Project Milestones & Impact",
            "content": [
                "Robust Autonomous Loop: Successfully moved from simple prompts to a structured 'Plan-then-Execute' workflow.",
                "Scalable Intelligence: Implemented a sub-agent architecture allowing the master agent to delegate complex sub-problems.",
                "Enterprise-Ready Infrastructure: Created a resilient networking layer with automatic retries and flexible configuration via .agentrc.",
                "Observability: Delivered a high-fidelity TUI that makes the 'black box' of AI reasoning transparent to the developer."
            ]
        }
    ]

    for item in slides_data:
        slide_layout = prs.slide_layouts[1] # Bullet layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = item["title"]
        
        # Set body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = item["content"][0]
        
        for bullet in item["content"][1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save('project_summary.pptx')
    print("Successfully generated project_summary.pptx")

if __name__ == "__main__":
    create_ppt()
